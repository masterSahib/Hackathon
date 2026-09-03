import os
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def build_presentation():
    input_pptx = r'C:\Users\a\Downloads\SIH2026-IDEA-Presentation-Format (1).pptx'
    output_pptx = r'C:\Users\a\OneDrive\Desktop\new_project\SIH2026_LabelTruth_Final_Submission.pptx'
    
    prs = pptx.Presentation(input_pptx)
    print(f"Loaded template with {len(prs.slides)} slides.")

    # Colors
    c_navy = RGBColor(15, 23, 42)       # #0F172A
    c_blue = RGBColor(2, 132, 199)      # #0284C7
    c_slate = RGBColor(51, 65, 85)      # #334155
    c_dark = RGBColor(30, 41, 59)       # #1E293B

    # ----------------------------------------------------
    # SLIDE 1: Title Slide
    # ----------------------------------------------------
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.shape_id == 93 and shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            
            p0 = tf.paragraphs[0]
            p0.text = "Problem Statement ID: SIH26034"
            p0.font.bold = True
            p0.font.size = Pt(14)
            p0.font.color.rgb = c_blue
            
            p1 = tf.add_paragraph()
            p1.text = "Problem Statement Title: Software System to check compliance of Packaged Commodities under Legal Metrology(Packaged Commodities) Rules, 2011 by scanning products, images and labels."
            p1.font.bold = True
            p1.font.size = Pt(11.5)
            p1.font.color.rgb = c_navy
            p1.space_before = Pt(4)
            
            p2 = tf.add_paragraph()
            p2.text = "Project / Solution Title: LabelTruth — AI-Powered Legal Metrology & Packaging Compliance Detector"
            p2.font.bold = True
            p2.font.size = Pt(12.5)
            p2.font.color.rgb = RGBColor(5, 150, 105) # Green
            p2.space_before = Pt(4)
            
            p3 = tf.add_paragraph()
            p3.text = "Theme: Agriculture, FoodTech & Rural Development | Category: Software"
            p3.font.size = Pt(11)
            p3.font.color.rgb = c_slate
            p3.space_before = Pt(3)
            
            p4 = tf.add_paragraph()
            p4.text = "Ministry / Dept: Department of Consumer Affairs (DoCA), Ministry of Consumer Affairs, Food & Public Distribution"
            p4.font.size = Pt(11)
            p4.font.color.rgb = c_slate
            p4.space_before = Pt(3)

            p5 = tf.add_paragraph()
            p5.text = "Team ID: [Your Team ID]  |  Team Name: [Your Team Name Registered on Portal]"
            p5.font.bold = True
            p5.font.size = Pt(11)
            p5.font.color.rgb = c_navy
            p5.space_before = Pt(4)

    # Helper function to add slide title if missing
    def set_slide_title(slide, title_text):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.shape_id in [111, 122, 134, 146]:
                shape.text_frame.text = title_text
                shape.text_frame.paragraphs[0].font.bold = True
                shape.text_frame.paragraphs[0].font.size = Pt(22)
                shape.text_frame.paragraphs[0].font.color.rgb = c_navy
                return shape
        # If slide 2 has no title shape, create one
        title_box = slide.shapes.add_textbox(Inches(0.66), Inches(0.2), Inches(10.0), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = c_navy
        return title_box

    # Helper function to clear existing placeholder content shapes (keeping footers/watermarks)
    def clean_and_get_content_box(slide, excluded_shape_ids):
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.shape_id in excluded_shape_ids:
                shapes_to_remove.append(shape)
        for s in shapes_to_remove:
            sp = s._element
            sp.getparent().remove(sp)
            
        # Add new structured content text box
        box = slide.shapes.add_textbox(Inches(0.66), Inches(1.1), Inches(12.0), Inches(5.15))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        return tf

    # ----------------------------------------------------
    # SLIDE 2: PROPOSED SOLUTION
    # ----------------------------------------------------
    slide2 = prs.slides[1]
    set_slide_title(slide2, "PROPOSED SOLUTION (Idea / Prototype)")
    tf2 = clean_and_get_content_box(slide2, [103]) # keep 100, 101, 102, 104
    
    # 1. Solution Overview
    p = tf2.paragraphs[0]
    p.text = "1. Detailed Explanation of Proposed Solution"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = c_blue
    
    p = tf2.add_paragraph()
    p.text = "• LabelTruth is an AI-powered compliance platform automating packaged commodity audits under Legal Metrology (Packaged Commodities) Rules, 2011 (LMPC Rules) and Legal Metrology Act, 2009 for Enforcement Officers (DoCA) and Consumers."
    p.font.size = Pt(10.5)
    p.font.color.rgb = c_navy

    p = tf2.add_paragraph()
    p.text = "• Dual-Pack Multimodal Vision AI (Gemini 2.5 Flash): Simultaneously extracts front/back text, Principal Display Panel (PDP) dimensions, barcodes, and mandatory declarations with zero manual entry."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf2.add_paragraph()
    p.text = "• Deterministic Rule Engine: Validates Rule 6(1)(a)-(h) [Mfg Details, Generic Commodity Name, Net Qty, Unit Sale Price (USP), MRP, Packing Date, Consumer Care] and Rules 7-9 [PDP Font Sizing]."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    # 2. How it addresses problem
    p = tf2.add_paragraph()
    p.text = "2. How it Addresses the Problem"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = c_blue
    p.space_before = Pt(6)

    p = tf2.add_paragraph()
    p.text = "• Eliminates Manual Bottlenecks: Replaces 15-minute physical audits with < 3-second automated verification across physical retail & e-commerce."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf2.add_paragraph()
    p.text = "• Zero-Tolerance Defect Detection: Detects missing USP, non-compliant font heights, and deceptive branding (e.g. Maida masked as Atta)."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf2.add_paragraph()
    p.text = "• Evidence Dossiers: Automatically generates court-admissible digital dockets with cropped photo proofs, timestamps, and violation codes."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    # 3. Innovation and Uniqueness
    p = tf2.add_paragraph()
    p.text = "3. Innovation and Uniqueness of the Solution"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = c_blue
    p.space_before = Pt(6)

    p = tf2.add_paragraph()
    p.text = "• Image-First Vision AI: Audits any unlisted or newly launched physical product directly without needing pre-existing barcode databases (unlike TruthIn)."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf2.add_paragraph()
    p.text = "• Interactive Regulatory AI Chatbot (/chat/product): Real-time conversational legal assistant for officers & citizens to query statutory thresholds."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf2.add_paragraph()
    p.text = "• Automated Legal Notice Generator: 1-click generation of digitally stamped PDF violation notices under Sec 18 & 36 of Legal Metrology Act, 2009."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    # ----------------------------------------------------
    # SLIDE 3: TECHNICAL APPROACH
    # ----------------------------------------------------
    slide3 = prs.slides[2]
    set_slide_title(slide3, "TECHNICAL APPROACH")
    tf3 = clean_and_get_content_box(slide3, [114])

    p = tf3.paragraphs[0]
    p.text = "1. End-to-End 4-Stage Architectural Pipeline"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = c_blue

    p = tf3.add_paragraph()
    p.text = " [1. Ingestion] Dual-Camera Snap / E-Commerce Image URL / Barcode Scan"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = c_navy

    p = tf3.add_paragraph()
    p.text = " ──► [2. Multimodal AI Extraction] Gemini 2.5 Flash Vision isolations of PDP area, OCR text & structured JSON declarations."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf3.add_paragraph()
    p.text = " ──► [3. Deterministic Compliance Engine] Mathematical verification of LMPC Rules 2011 (Rule 6, 7-9), USP formulas, and FSSAI rules."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf3.add_paragraph()
    p.text = " ──► [4. Enforcement & Reporting] 100-Point Truth Score dashboard, ReportLab Statutory PDF Notices, and Interactive AI Chat."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf3.add_paragraph()
    p.text = "2. Technology Stack & Frameworks"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = c_blue
    p.space_before = Pt(8)

    p = tf3.add_paragraph()
    p.text = "• Mobile & Web Frontend: Flutter (Dart) with Riverpod, CameraX, and Text-to-Speech (TTS) for high-performance cross-platform execution."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf3.add_paragraph()
    p.text = "• Backend API: FastAPI (Python 3.10) with asynchronous non-blocking routes (<100ms response time) and Pydantic v2 data validation."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf3.add_paragraph()
    p.text = "• Database & Cloud: PostgreSQL (Render Cloud) with SQLAlchemy ORM for encrypted inspection logs, evidence photos, and audit histories."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf3.add_paragraph()
    p.text = "3. Key Engineering Methodologies"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = c_blue
    p.space_before = Pt(8)

    p = tf3.add_paragraph()
    p.text = "• Zero-Hallucination Separation: Multimodal AI is strictly used for text extraction; all legal evaluations are 100% computed mathematically."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf3.add_paragraph()
    p.text = "• Hybrid Decoupled Scanning: Operates via Dual-Camera packaging OCR or Open Food Facts live barcode registry with automatic fallback."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    # ----------------------------------------------------
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ----------------------------------------------------
    slide4 = prs.slides[3]
    set_slide_title(slide4, "FEASIBILITY AND VIABILITY")
    tf4 = clean_and_get_content_box(slide4, [123, 126])

    p = tf4.paragraphs[0]
    p.text = "1. Feasibility Analysis"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = c_blue

    p = tf4.add_paragraph()
    p.text = "• Technical Feasibility: Fully functional working prototype deployed and tested (Flutter + FastAPI + PostgreSQL + Gemini 2.5 Flash). Delivers >95% OCR accuracy and <3s audit response on real Indian retail products."
    p.font.size = Pt(10)
    p.font.color.rgb = c_navy

    p = tf4.add_paragraph()
    p.text = "• Operational Feasibility: Zero learning curve for field officers; runs seamlessly on standard smartphones without costly dedicated hardware."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf4.add_paragraph()
    p.text = "• Legal & Financial Viability: Strictly aligned with Legal Metrology Act, 2009 & LMPC Rules, 2011. Negligible operational cost (< ₹0.15/scan)."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf4.add_paragraph()
    p.text = "2. Potential Challenges, Risks & Mitigation Strategies"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = c_blue
    p.space_before = Pt(8)

    p = tf4.add_paragraph()
    p.text = "• Challenge 1: Packaging Glare, Crinkling & Curved Surfaces ──► Strategy: Live camera bounding boxes, auto-flash assist & multi-angle sharpness guidance."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf4.add_paragraph()
    p.text = "• Challenge 2: Multilingual Packaging Across Indian States ──► Strategy: Gemini multimodal vision inherently parses 12+ Indian languages (Hindi, Tamil, Marathi, etc.)."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf4.add_paragraph()
    p.text = "• Challenge 3: AI Hallucinations in Legal Notices ──► Strategy: Deterministic rule engine strictly separates perceptual text from statutory logic."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf4.add_paragraph()
    p.text = "• Challenge 4: Low Connectivity in Rural Markets ──► Strategy: Offline-first local SQLite cache queues inspections and auto-syncs when online."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    # ----------------------------------------------------
    # SLIDE 5: IMPACT AND BENEFITS
    # ----------------------------------------------------
    slide5 = prs.slides[4]
    set_slide_title(slide5, "IMPACT AND BENEFITS")
    tf5 = clean_and_get_content_box(slide5, [135, 138])

    p = tf5.paragraphs[0]
    p.text = "1. Multi-Stakeholder Impact"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = c_blue

    p = tf5.add_paragraph()
    p.text = "• Department of Consumer Affairs (DoCA) & Enforcement Officers: 10x faster market surveillance, automated evidence collection, and digital court dockets."
    p.font.size = Pt(10)
    p.font.color.rgb = c_navy

    p = tf5.add_paragraph()
    p.text = "• 1.4 Billion Indian Consumers: Complete price transparency (accurate MRP & Unit Sale Price), elimination of deceptive packaging, and 1-click NCH reporting."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf5.add_paragraph()
    p.text = "• FMCG Manufacturers & Retailers: Pre-market digital label verification prevents crores in packaging recall losses and regulatory penalties."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf5.add_paragraph()
    p.text = "2. Core Value Dimensions (Social, Economic, Environmental)"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = c_blue
    p.space_before = Pt(8)

    p = tf5.add_paragraph()
    p.text = "• Governance & Legal: Standardizes automated enforcement under Sections 18 & 36 of Legal Metrology Act, 2009 with zero subjective human bias."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf5.add_paragraph()
    p.text = "• Economic & Fair Trade: Eliminates dual-MRP and hidden price inflation by enforcing transparent Unit Sale Pricing (USP per g/ml)."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf5.add_paragraph()
    p.text = "• Social & Inclusivity: Text-to-Speech (TTS) voice readout empowers rural, elderly, and visually impaired citizens to verify packaging integrity."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf5.add_paragraph()
    p.text = "• Environmental Sustainability: Pre-market label verification prevents millions of non-compliant printed plastic/paper packages from being scrapped."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    # ----------------------------------------------------
    # SLIDE 6: RESEARCH AND REFERENCES
    # ----------------------------------------------------
    slide6 = prs.slides[5]
    set_slide_title(slide6, "RESEARCH AND REFERENCES")
    tf6 = clean_and_get_content_box(slide6, [147, 150])

    p = tf6.paragraphs[0]
    p.text = "1. Statutory Acts, Rules & Government Portals"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = c_blue

    p = tf6.add_paragraph()
    p.text = "• Legal Metrology Act, 2009 (Act 1 of 2010): Sections 18 & 36 (Statutory standards & penal liability) ── consumeraffairs.gov.in/pages/legal-metrology-act"
    p.font.size = Pt(10)
    p.font.color.rgb = c_navy

    p = tf6.add_paragraph()
    p.text = "• Legal Metrology (Packaged Commodities) Rules, 2011: Rule 6 (Mandatory Declarations & USP), Rules 7-9 (PDP Font Sizing) ── consumeraffairs.gov.in"
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf6.add_paragraph()
    p.text = "• FSSAI Labelling & Display Regulations 2020 & Claims 2018: Sec 23 (Descending Hierarchy) & HFSS Standards ── fssai.gov.in"
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf6.add_paragraph()
    p.text = "• Consumer Protection Act, 2019 & CCPA Guidelines 2022: Unfair trade practices & misleading advertising bans ── consumerhelpline.gov.in"
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf6.add_paragraph()
    p.text = "2. Technical Literature & Open Data Standards"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = c_blue
    p.space_before = Pt(8)

    p = tf6.add_paragraph()
    p.text = "• Multimodal Large Language Models for Scene Text & Document Understanding (Google DeepMind / Gemini Vision Research)."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf6.add_paragraph()
    p.text = "• Hybrid Neuro-Symbolic Compliance Verification: Eliminating LLM hallucinations via deterministic mathematical rule separation."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    p = tf6.add_paragraph()
    p.text = "• Open Food Facts India & Global Packaging Schema (world.openfoodfacts.org) for real-time barcode product registry validation."
    p.font.size = Pt(10)
    p.font.color.rgb = c_slate

    # Delete slide 7 (Instructions slide) as per note on slide 7
    if len(prs.slides) > 6:
        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]
        print("Removed Slide 7 (Instructions slide) to comply with 6-slide max limit.")

    prs.save(output_pptx)
    print(f"Successfully generated final presentation at: {output_pptx}")

    # Also export to PDF
    try:
        import win32com.client
        powerpoint = win32com.client.Dispatch('PowerPoint.Application')
        ppt_abs = os.path.abspath(output_pptx)
        pdf_abs = os.path.abspath(r'C:\Users\a\OneDrive\Desktop\new_project\SIH2026_LabelTruth_Final_Submission.pdf')
        deck = powerpoint.Presentations.Open(ppt_abs, WithWindow=False)
        deck.SaveAs(pdf_abs, 32)
        deck.Close()
        powerpoint.Quit()
        print(f"Successfully exported final PDF at: {pdf_abs}")
    except Exception as e:
        print("PDF export note:", e)

if __name__ == "__main__":
    build_presentation()
